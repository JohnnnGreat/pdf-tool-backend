"""Document service — Word, Excel, PowerPoint operations."""
import io
import os
import zipfile
from pathlib import Path
from typing import Optional

from fastapi import HTTPException


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------

def create_docx(data: dict, output_path: str) -> None:
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    if title := data.get("title"):
        doc.add_heading(title, level=1)
    for section in data.get("sections", []):
        if heading := section.get("heading"):
            doc.add_heading(heading, level=section.get("level", 2))
        if text := section.get("text"):
            doc.add_paragraph(text)
    doc.save(output_path)


def find_replace_docx(input_path: str, output_path: str, find: str, replace: str) -> int:
    from docx import Document
    doc = Document(input_path)
    count = 0
    for para in doc.paragraphs:
        if find in para.text:
            for run in para.runs:
                if find in run.text:
                    run.text = run.text.replace(find, replace)
                    count += 1
    doc.save(output_path)
    return count


def docx_to_markdown(input_path: str, output_path: str) -> None:
    from docx import Document
    doc = Document(input_path)
    lines = []
    for para in doc.paragraphs:
        style = para.style.name.lower()
        if "heading 1" in style:
            lines.append(f"# {para.text}")
        elif "heading 2" in style:
            lines.append(f"## {para.text}")
        elif "heading 3" in style:
            lines.append(f"### {para.text}")
        else:
            lines.append(para.text)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n\n".join(lines))


def merge_docx_files(input_paths: list[str], output_path: str) -> None:
    from docx import Document
    from docx.oxml.ns import qn
    merged = Document(input_paths[0])
    for path in input_paths[1:]:
        doc = Document(path)
        merged.add_page_break()
        for element in doc.element.body:
            merged.element.body.append(element)
    merged.save(output_path)


def compare_docx_files(path_a: str, path_b: str) -> dict:
    from docx import Document
    doc_a = Document(path_a)
    doc_b = Document(path_b)
    paras_a = [p.text for p in doc_a.paragraphs]
    paras_b = [p.text for p in doc_b.paragraphs]
    added = [p for p in paras_b if p not in paras_a]
    removed = [p for p in paras_a if p not in paras_b]
    return {"added": added, "removed": removed, "total_changes": len(added) + len(removed)}


def extract_text_from_file(input_path: str) -> str:
    ext = Path(input_path).suffix.lower()
    if ext == ".pdf":
        import fitz
        doc = fitz.open(input_path)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text
    elif ext in (".doc", ".docx"):
        from docx import Document
        doc = Document(input_path)
        return "\n".join(p.text for p in doc.paragraphs)
    elif ext in (".ppt", ".pptx"):
        from pptx import Presentation
        prs = Presentation(input_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    raise HTTPException(status_code=415, detail=f"Unsupported format: {ext}")


def mail_merge(template_path: str, csv_path: str, output_dir: str) -> list[str]:
    import csv
    from docx import Document
    paths = []
    with open(csv_path, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for i, row in enumerate(reader):
            doc = Document(template_path)
            for para in doc.paragraphs:
                for key, value in row.items():
                    for run in para.runs:
                        run.text = run.text.replace(f"{{{{{key}}}}}", value)
            out = os.path.join(output_dir, f"document_{i + 1}.docx")
            doc.save(out)
            paths.append(out)
    return paths


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------

def merge_pptx_files(input_paths: list[str], output_path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    merged = Presentation(input_paths[0])
    for path in input_paths[1:]:
        prs = Presentation(path)
        for slide in prs.slides:
            xml = slide._element
            merged.slides._sldIdLst.append(xml)
    merged.save(output_path)


def pptx_to_images(input_path: str, output_dir: str, fmt: str = "png", dpi: int = 150) -> list[str]:
    import fitz
    doc = fitz.open(input_path)
    paths = []
    for i, page in enumerate(doc):
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        path = os.path.join(output_dir, f"slide_{i + 1}.{fmt}")
        pix.save(path)
        paths.append(path)
    doc.close()
    return paths


def extract_slide_notes(input_path: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(input_path)
    notes = []
    for i, slide in enumerate(prs.slides):
        note_text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            note_text = tf.text if tf else ""
        notes.append({"slide": i + 1, "notes": note_text})
    return notes


def images_to_pptx(input_paths: list[str], output_path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for img_path in input_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    prs.save(output_path)


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------

def merge_excel_files(input_paths: list[str], output_path: str) -> None:
    from openpyxl import load_workbook, Workbook
    wb_out = Workbook()
    wb_out.remove(wb_out.active)
    for path in input_paths:
        wb = load_workbook(path)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            new_ws = wb_out.create_sheet(title=f"{Path(path).stem}_{sheet_name}")
            for row in ws.iter_rows(values_only=True):
                new_ws.append(list(row))
    wb_out.save(output_path)


def split_excel_sheets(input_path: str, output_dir: str) -> list[str]:
    from openpyxl import load_workbook, Workbook
    wb = load_workbook(input_path)
    paths = []
    for sheet_name in wb.sheetnames:
        new_wb = Workbook()
        ws = new_wb.active
        ws.title = sheet_name
        for row in wb[sheet_name].iter_rows(values_only=True):
            ws.append(list(row))
        out = os.path.join(output_dir, f"{sheet_name}.xlsx")
        new_wb.save(out)
        paths.append(out)
    return paths


def clean_excel(input_path: str, output_path: str) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(input_path)
    removed_rows = 0
    for ws in wb.worksheets:
        blank_rows = [row[0].row for row in ws.iter_rows() if all(cell.value is None for cell in row)]
        for r in reversed(blank_rows):
            ws.delete_rows(r)
            removed_rows += 1
    wb.save(output_path)
    return {"blank_rows_removed": removed_rows}


def json_to_excel(input_path: str, output_path: str) -> None:
    import json
    from openpyxl import Workbook
    with open(input_path, encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, list):
        data = [data]
    wb = Workbook()
    ws = wb.active
    if data:
        ws.append(list(data[0].keys()))
        for row in data:
            ws.append(list(row.values()))
    wb.save(output_path)


def excel_to_json(input_path: str, output_path: str) -> None:
    import json
    from openpyxl import load_workbook
    wb = load_workbook(input_path, read_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        data = []
    else:
        headers = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(rows[0])]
        data = [dict(zip(headers, row)) for row in rows[1:]]
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, default=str)
