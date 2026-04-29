"""Format conversion service — PDF ↔ Word/Excel/PPT/Image/HTML/Markdown/CSV/Text."""
import logging
import os
import subprocess
from pathlib import Path

from fastapi import HTTPException

from app.core.config import settings
from rust_converter import (
    DocxGenerationError as RustDocxGenerationError,
    InvalidPdfError as RustInvalidPdfError,
    RustConversionError,
    RustModuleNotBuiltError,
    UnsupportedScannedPdfError as RustUnsupportedScannedPdfError,
    convert_pdf_to_docx as rust_convert_pdf_to_docx,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# LibreOffice helper (Word/Excel/PPTX → PDF, EPUB → PDF, SVG → PNG)
# ---------------------------------------------------------------------------

def _resolve_libreoffice_path() -> str:
    raw = settings.LIBREOFFICE_PATH or "soffice"
    raw = raw.strip().strip('"')
    if raw == "soffice":
        return raw
    resolved = Path(raw).resolve()
    if not resolved.exists():
        raise HTTPException(status_code=500, detail=f"LibreOffice binary not found at '{resolved}'")
    return str(resolved)


def _libreoffice_convert(
    input_path: str,
    output_dir: str,
    target_format: str,
    extra_args: list[str] | None = None,
) -> str:
    lo = _resolve_libreoffice_path()
    cmd = [lo, "--headless"]
    if extra_args:
        cmd.extend(extra_args)
    cmd += ["--convert-to", target_format, "--outdir", output_dir, input_path]
    # Set HOME to a writable temp dir so LibreOffice can create its user profile.
    # Without this, non-root Docker users see "no export filter found" errors.
    env = os.environ.copy()
    env.setdefault("HOME", "/tmp")
    result = subprocess.run(
        cmd,
        capture_output=True, text=True, timeout=120, env=env,
    )
    if result.returncode != 0:
        logger.debug("LO stdout: %s", result.stdout)
        logger.debug("LO stderr: %s", result.stderr)
        raise HTTPException(status_code=500,
                            detail=f"LibreOffice conversion failed: {result.stderr or result.stdout}")
    # LibreOffice names the output after the input stem; find the actual file.
    ext = target_format.split(":")[0].lower()   # handles "docx:MS Word 2007 XML" → "docx"
    stem = Path(input_path).stem
    matches = [
        f for f in os.listdir(output_dir)
        if f.lower().startswith(stem.lower()) and f.lower().endswith(f".{ext}")
    ]
    if not matches:
        raise HTTPException(status_code=500,
                            detail=f"Conversion output not found for {stem}.{ext}")
    return os.path.join(output_dir, matches[0])


# ---------------------------------------------------------------------------
# Office → PDF  (LibreOffice is gold standard for these)
# ---------------------------------------------------------------------------

def word_to_pdf(input_path: str, output_dir: str) -> str:
    return _libreoffice_convert(input_path, output_dir, "pdf")


def excel_to_pdf(input_path: str, output_dir: str) -> str:
    return _libreoffice_convert(input_path, output_dir, "pdf")


def pptx_to_pdf(input_path: str, output_dir: str) -> str:
    return _libreoffice_convert(input_path, output_dir, "pdf")


# ---------------------------------------------------------------------------
# PDF → Word  — multi-strategy engine
# ---------------------------------------------------------------------------
# Strategy selection:
#   1. Analyse: count extractable chars to decide digital vs scanned
#   2a. Digital → pdf2docx  (best-in-class layout reconstruction)
#         If pdf2docx fails → custom PyMuPDF + python-docx extractor
#         If that fails     → LibreOffice writer_pdf_import (emergency)
#   2b. Scanned → Tesseract per-page OCR → structured python-docx
# ---------------------------------------------------------------------------

def _count_chars(input_path: str, pages: int = 5) -> int:
    """Return the number of extractable text chars from the first N pages."""
    import fitz
    doc = fitz.open(input_path)
    total = 0
    for i in range(min(pages, len(doc))):
        total += len(doc[i].get_text().strip())
    doc.close()
    return total


def _pdf_to_word_pdf2docx(input_path: str, output_path: str) -> None:
    """Strategy 1: pdf2docx — best layout fidelity for digital PDFs."""
    from pdf2docx import Converter
    import fitz

    # Pre-check: count pages and chars so we can validate output quality
    doc_check = fitz.open(input_path)
    num_pages = len(doc_check)
    total_chars = sum(len(doc_check[i].get_text().strip()) for i in range(min(num_pages, 5)))
    doc_check.close()

    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()

    # pdf2docx sometimes writes a near-empty file on silent failure.
    # Use a size threshold proportional to content: at least 1 KB per page
    # or 2 KB minimum, whichever is larger.
    min_expected = max(2048, num_pages * 1024)
    if not os.path.exists(output_path) or os.path.getsize(output_path) < min_expected:
        raise RuntimeError(
            f"pdf2docx produced a suspiciously small output "
            f"({os.path.getsize(output_path) if os.path.exists(output_path) else 0} bytes "
            f"for {num_pages} pages, expected >{min_expected})"
        )

    # Secondary quality check: open the docx and verify it has meaningful text
    try:
        from docx import Document as _Doc
        _d = _Doc(output_path)
        extracted = " ".join(p.text for p in _d.paragraphs).strip()
        # If the PDF had text but the docx has almost none, it's a bad conversion
        if total_chars > 200 and len(extracted) < max(50, total_chars * 0.10):
            raise RuntimeError(
                f"pdf2docx output has too little text ({len(extracted)} chars) "
                f"compared to source PDF ({total_chars} chars)"
            )
    except RuntimeError:
        raise
    except Exception as e:
        logger.debug("pdf2docx quality check skipped: %s", e)


def _pdf_to_word_pymupdf(input_path: str, output_path: str) -> None:
    """
    Strategy 2: Custom PyMuPDF + python-docx extractor.

    Improvements over the original:
    - Groups lines into paragraphs by vertical proximity (instead of one para per line)
    - Preserves text alignment (left / center / right / justify) per block
    - Preserves font family from PDF spans (not just Calibri)
    - Detects headings by font size AND bold flag
    - Embeds images at correct proportional width
    - Extracts tables via pdfplumber with header styling
    - Sets page size to match the PDF page dimensions
    """
    import fitz
    from docx import Document
    from docx.shared import Pt, Inches, Emu, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import io

    try:
        import pdfplumber
        HAS_PLUMBER = True
    except ImportError:
        HAS_PLUMBER = False

    doc_out = Document()

    # ── Style defaults ────────────────────────────────────────────────────────
    style = doc_out.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    pdf = fitz.open(input_path)

    # Pre-compute per-page table bounding boxes so we can skip those blocks
    page_tables: dict[int, list] = {}
    if HAS_PLUMBER:
        try:
            with pdfplumber.open(input_path) as plumb:
                for p_idx, page in enumerate(plumb.pages):
                    tbls = page.find_tables()
                    page_tables[p_idx] = [t.bbox for t in tbls] if tbls else []
        except Exception:
            pass

    def _in_table(bbox, p_idx: int) -> bool:
        x0, y0, x1, y1 = bbox
        for tb in page_tables.get(p_idx, []):
            tx0, ty0, tx1, ty1 = tb
            if x0 >= tx0 - 2 and y0 >= ty0 - 2 and x1 <= tx1 + 2 and y1 <= ty1 + 2:
                return True
        return False

    def _add_table_from_plumber(p_idx: int) -> None:
        if not HAS_PLUMBER:
            return
        try:
            with pdfplumber.open(input_path) as plumb:
                page = plumb.pages[p_idx]
                for tbl in page.extract_tables() or []:
                    if not tbl:
                        continue
                    rows = [[str(c or "") for c in row] for row in tbl]
                    t = doc_out.add_table(rows=len(rows), cols=len(rows[0]))
                    t.style = "Table Grid"
                    for r_i, row in enumerate(rows):
                        for c_i, val in enumerate(row):
                            cell = t.cell(r_i, c_i)
                            cell.text = val
                            if r_i == 0:
                                for run in cell.paragraphs[0].runs:
                                    run.bold = True
        except Exception as e:
            logger.debug("pdfplumber table extract failed: %s", e)

    def _get_alignment(block, page_width: float) -> WD_ALIGN_PARAGRAPH:
        """Infer paragraph alignment from block position on the page."""
        x0, _, x1, _ = block["bbox"]
        block_center = (x0 + x1) / 2
        page_center = page_width / 2
        block_width = x1 - x0

        # Center-aligned: block center is near page center and block is not full-width
        if abs(block_center - page_center) < page_width * 0.08 and block_width < page_width * 0.7:
            return WD_ALIGN_PARAGRAPH.CENTER
        # Right-aligned: block starts in the right 40% of the page
        if x0 > page_width * 0.6:
            return WD_ALIGN_PARAGRAPH.RIGHT
        return WD_ALIGN_PARAGRAPH.LEFT

    def _group_lines_into_paragraphs(lines: list) -> list[list]:
        """
        Group consecutive lines into paragraphs based on vertical gap.
        A gap larger than 1.5× the line height signals a new paragraph.
        """
        if not lines:
            return []
        groups = [[lines[0]]]
        for line in lines[1:]:
            prev = groups[-1][-1]
            prev_y1 = prev["bbox"][3]
            curr_y0 = line["bbox"][1]
            # Estimate line height from the previous line's spans
            prev_heights = [s["size"] for s in prev["spans"] if s["text"].strip()]
            line_h = max(prev_heights) if prev_heights else 12
            gap = curr_y0 - prev_y1
            if gap > line_h * 1.2:
                groups.append([line])
            else:
                groups[-1].append(line)
        return groups

    def _safe_font(font_name: str) -> str:
        """Map PDF font names to safe Word-compatible equivalents."""
        fn = font_name.lower()
        if any(x in fn for x in ("arial", "helvetica", "sans")):
            return "Arial"
        if any(x in fn for x in ("times", "serif", "roman")):
            return "Times New Roman"
        if any(x in fn for x in ("courier", "mono", "consol")):
            return "Courier New"
        if "georgia" in fn:
            return "Georgia"
        if "verdana" in fn:
            return "Verdana"
        if "calibri" in fn:
            return "Calibri"
        # Strip subset prefix like "ABCDEF+FontName"
        if "+" in font_name:
            base = font_name.split("+", 1)[1]
            return _safe_font(base)
        return "Calibri"  # safe default

    for p_idx, page in enumerate(pdf):
        page_width = page.rect.width
        page_height = page.rect.height

        # Set the section page size to match the PDF page (first page only)
        if p_idx == 0:
            section = doc_out.sections[0]
            # PDF points → EMU (1 pt = 12700 EMU)
            section.page_width  = Emu(int(page_width  * 12700))
            section.page_height = Emu(int(page_height * 12700))
            # Proportional margins: ~1 inch on each side
            margin = Emu(int(72 * 12700))  # 72 pt = 1 inch
            section.left_margin   = margin
            section.right_margin  = margin
            section.top_margin    = margin
            section.bottom_margin = margin

        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

        # Estimate body font size for heading detection (median of all span sizes)
        sizes = []
        for b in blocks:
            if b["type"] != 0:
                continue
            for line in b["lines"]:
                for span in line["spans"]:
                    if span["text"].strip():
                        sizes.append(round(span["size"]))
        body_size = sorted(sizes)[len(sizes) // 2] if sizes else 11

        tables_added_for_page = False

        for b in blocks:
            if b["type"] == 1:
                # Image block — embed inline at proportional width
                try:
                    bx0, by0, bx1, by1 = b["bbox"]
                    img_w_pt = bx1 - bx0
                    # Convert PDF points to inches (72 pt/inch), cap at usable width
                    usable_w = (page_width - 144) / 72  # page minus 2-inch margins
                    img_w_in = min(img_w_pt / 72, usable_w)
                    img_bytes = page.get_pixmap(clip=fitz.Rect(b["bbox"]), dpi=150).tobytes("png")
                    buf = io.BytesIO(img_bytes)
                    doc_out.add_picture(buf, width=Inches(img_w_in))
                except Exception:
                    pass
                continue

            if b["type"] != 0:
                continue

            # Skip blocks that fall inside a detected table
            if _in_table(b["bbox"], p_idx):
                if not tables_added_for_page:
                    _add_table_from_plumber(p_idx)
                    tables_added_for_page = True
                continue

            alignment = _get_alignment(b, page_width)
            all_lines = b["lines"]
            para_groups = _group_lines_into_paragraphs(all_lines)

            for group in para_groups:
                # Collect full text and dominant properties for this paragraph group
                group_text = " ".join(
                    "".join(s["text"] for s in line["spans"]).strip()
                    for line in group
                ).strip()

                if not group_text:
                    continue

                # Determine heading level from font size delta vs body
                max_size = max(
                    (round(s["size"]) for line in group for s in line["spans"] if s["text"].strip()),
                    default=body_size,
                )
                delta = max_size - body_size

                # Also check if all spans in the group are bold (another heading signal)
                all_bold = all(
                    bool(s["flags"] & 2**4)
                    for line in group for s in line["spans"] if s["text"].strip()
                )

                if delta >= 8 or (delta >= 4 and all_bold):
                    para = doc_out.add_heading(group_text, level=1)
                    para.alignment = alignment
                elif delta >= 5:
                    para = doc_out.add_heading(group_text, level=2)
                    para.alignment = alignment
                elif delta >= 3:
                    para = doc_out.add_heading(group_text, level=3)
                    para.alignment = alignment
                else:
                    para = doc_out.add_paragraph()
                    para.alignment = alignment
                    # Add runs span-by-span to preserve inline formatting
                    for line in group:
                        for span in line["spans"]:
                            if not span["text"]:
                                continue
                            run = para.add_run(span["text"])
                            run.bold   = bool(span["flags"] & 2**4)
                            run.italic = bool(span["flags"] & 2**1)
                            run.font.size = Pt(round(span["size"]))
                            # Font family
                            font_name = _safe_font(span.get("font", ""))
                            run.font.name = font_name
                            # Colour
                            c = span.get("color", 0)
                            r, g, bv = (c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF
                            if (r, g, bv) != (0, 0, 0):
                                run.font.color.rgb = RGBColor(r, g, bv)
                        # Add a soft line break between lines within the same paragraph group
                        # (except after the last line)
                        if line is not group[-1]:
                            run = para.add_run()
                            run.add_break()

        # Page break between pages (except last)
        if p_idx < len(pdf) - 1:
            doc_out.add_page_break()

    pdf.close()
    doc_out.save(output_path)

    if os.path.getsize(output_path) < 512:
        raise RuntimeError("PyMuPDF extractor produced an empty output")


def _pdf_to_word_scanned(input_path: str, output_path: str) -> None:
    """
    Strategy for scanned PDFs: render each page → Tesseract OCR →
    build a clean DOCX preserving paragraph structure.
    """
    import fitz
    import io
    from docx import Document
    from docx.shared import Pt, Inches

    try:
        import pytesseract
        from PIL import Image
        if settings.TESSERACT_PATH:
            pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_PATH
    except ImportError:
        # No Tesseract → fall back to digital extractor anyway
        _pdf_to_word_pymupdf(input_path, output_path)
        return

    doc_out = Document()
    for section in doc_out.sections:
        section.left_margin = section.right_margin = Inches(1.0)
        section.top_margin  = section.bottom_margin = Inches(1.0)

    pdf = fitz.open(input_path)
    for p_idx, page in enumerate(pdf):
        mat = fitz.Matrix(2.0, 2.0)          # 2× = 144 DPI — good OCR quality
        pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        raw_text = pytesseract.image_to_string(img, config="--psm 1 --oem 3")

        for line in raw_text.splitlines():
            stripped = line.strip()
            if not stripped:
                doc_out.add_paragraph()     # preserve blank lines
                continue
            para = doc_out.add_paragraph(stripped)
            para.runs[0].font.size = Pt(11)
            para.runs[0].font.name = "Calibri"

        if p_idx < len(pdf) - 1:
            doc_out.add_page_break()

    pdf.close()
    doc_out.save(output_path)


def _pdf_to_word_libreoffice(input_path: str, output_path: str) -> None:
    """Emergency fallback: LibreOffice writer_pdf_import."""
    output_dir = os.path.dirname(output_path)
    lo_out = _libreoffice_convert(
        input_path,
        output_dir,
        "docx:MS Word 2007 XML",
        extra_args=["--infilter=writer_pdf_import"],
    )
    if lo_out != output_path and os.path.exists(lo_out):
        os.replace(lo_out, output_path)


def _legacy_pdf_to_word(input_path: str, output_path: str) -> None:
    """
    Smart PDF → DOCX conversion with cascading strategy selection.

    Digital PDFs  → pdf2docx (best layout) → PyMuPDF+python-docx → LibreOffice
    Scanned PDFs  → Tesseract OCR → structured DOCX
    """
    char_count = _count_chars(input_path)
    is_scanned = char_count < 80

    logger.info("pdf_to_word: chars=%d scanned=%s path=%s", char_count, is_scanned, input_path)

    if is_scanned:
        try:
            _pdf_to_word_scanned(input_path, output_path)
            logger.info("pdf_to_word: scanned strategy succeeded")
            return
        except Exception as e:
            logger.warning("pdf_to_word: scanned strategy failed (%s) — falling back to LibreOffice", e)
            _pdf_to_word_libreoffice(input_path, output_path)
            return

    # Digital PDF — try in order of quality
    for name, strategy in [
        ("pdf2docx",  _pdf_to_word_pdf2docx),
        ("PyMuPDF",   _pdf_to_word_pymupdf),
        ("LibreOffice", _pdf_to_word_libreoffice),
    ]:
        try:
            strategy(input_path, output_path)
            logger.info("pdf_to_word: strategy '%s' succeeded", name)
            return
        except Exception as e:
            logger.warning("pdf_to_word: strategy '%s' failed (%s) — trying next", name, e)
            # Clean up partial output before next attempt
            if os.path.exists(output_path):
                os.remove(output_path)

    raise HTTPException(
        status_code=500,
        detail="PDF to Word conversion failed. The file may be encrypted, corrupted, or in an unsupported format.",
    )



# ---------------------------------------------------------------------------
# PDF → Excel  (pdfplumber — best table extraction accuracy)
# ---------------------------------------------------------------------------

def pdf_to_word(input_path: str, output_path: str) -> None:
    """
    Convert a PDF to DOCX using the custom Rust visual replica engine.

    The Rust module extracts page sizes and text layout, renders each page
    to a background image for fidelity, then writes a DOCX with absolutely
    positioned text boxes over the rendered page image.
    """
    if os.path.exists(output_path):
        os.remove(output_path)

    try:
        success = rust_convert_pdf_to_docx(input_path, output_path)
    except RustModuleNotBuiltError as exc:
        logger.exception("Rust PDF-to-DOCX module is unavailable: %s", exc)
        raise HTTPException(
            status_code=500,
            detail=(
                "Rust PDF-to-DOCX converter is not installed. "
                "Build it with `maturin develop --release` from the backend root."
            ),
        ) from exc
    except RustInvalidPdfError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except RustUnsupportedScannedPdfError as exc:
        raise HTTPException(status_code=422, detail=str(exc)) from exc
    except (RustDocxGenerationError, RustConversionError) as exc:
        logger.warning("Rust PDF-to-DOCX conversion failed: %s", exc)
        raise HTTPException(status_code=500, detail=str(exc)) from exc
    except Exception as exc:  # pragma: no cover - defensive boundary
        logger.exception("Unexpected Rust PDF-to-DOCX failure: %s", exc)
        raise HTTPException(
            status_code=500,
            detail="PDF to Word conversion failed inside the Rust converter.",
        ) from exc

    if not success:
        raise HTTPException(
            status_code=500,
            detail="Rust PDF-to-DOCX converter returned an unsuccessful result.",
        )

    if not os.path.exists(output_path):
        raise HTTPException(
            status_code=500,
            detail="Rust PDF-to-DOCX converter did not create an output file.",
        )


def pdf_to_excel(input_path: str, output_path: str) -> None:
    """
    Extract tables from PDF into .xlsx.
    Uses pdfplumber with multiple strategies (lines → text fallback).
    Each detected table gets its own worksheet.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    BLUE   = "2563EB"
    STRIPE = "F1F5F9"
    BORDER_COLOR = "E2E8F0"

    def _border():
        s = Side(style="thin", color=BORDER_COLOR)
        return Border(left=s, right=s, top=s, bottom=s)

    try:
        import pdfplumber
        _extractor = "pdfplumber"
    except ImportError:
        _extractor = "pymupdf"

    wb = Workbook()
    wb.remove(wb.active)
    table_count = 0

    def _write_table(rows: list[list], sheet_title: str) -> None:
        nonlocal table_count
        table_count += 1
        ws = wb.create_sheet(title=sheet_title[:31])
        for r_idx, row in enumerate(rows, 1):
            for c_idx, val in enumerate(row, 1):
                cell = ws.cell(row=r_idx, column=c_idx, value=val or "")
                cell.border = _border()
                cell.alignment = Alignment(vertical="top", wrap_text=True)
                if r_idx == 1:
                    cell.fill = PatternFill("solid", fgColor=BLUE)
                    cell.font = Font(bold=True, color="FFFFFF", size=10)
                elif r_idx % 2 == 0:
                    cell.fill = PatternFill("solid", fgColor=STRIPE)
        # Auto column widths
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=8)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 60)
        ws.freeze_panes = "A2"

    if _extractor == "pdfplumber":
        import pdfplumber
        with pdfplumber.open(input_path) as pdf:
            for p_num, page in enumerate(pdf.pages, 1):
                # Try line-based strategy first, fall back to text-based
                for strategy in (
                    {"vertical_strategy": "lines",    "horizontal_strategy": "lines",
                     "intersection_tolerance": 5},
                    {"vertical_strategy": "text",     "horizontal_strategy": "text",
                     "snap_tolerance": 3},
                ):
                    tables = page.extract_tables(strategy)
                    if tables:
                        break
                for t_idx, tbl in enumerate(tables or [], 1):
                    if tbl:
                        _write_table(tbl, f"P{p_num}-T{t_idx}")
    else:
        import fitz
        doc = fitz.open(input_path)
        for p_num, page in enumerate(doc, 1):
            for t_idx, tbl in enumerate(page.find_tables().tables, 1):
                rows = tbl.extract()
                if rows:
                    _write_table(rows, f"P{p_num}-T{t_idx}")
        doc.close()

    if not wb.worksheets:
        ws = wb.create_sheet("No Tables Found")
        ws.cell(1, 1, "No tables were detected in this PDF.")
        ws.column_dimensions["A"].width = 40

    wb.save(output_path)


# ---------------------------------------------------------------------------
# PDF → PPTX  (visual fidelity: each page → full-bleed slide image)
# ---------------------------------------------------------------------------

def pdf_to_pptx(input_path: str, output_dir: str) -> str:
    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    doc = fitz.open(input_path)
    prs = Presentation()

    if doc:
        p0 = doc[0]
        # Convert PDF points (72 pt/inch) → EMU (914400 EMU/inch)
        prs.slide_width  = Emu(int(p0.rect.width  / 72 * 914400))
        prs.slide_height = Emu(int(p0.rect.height / 72 * 914400))

    blank_layout = prs.slide_layouts[6]

    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=200)   # 200 DPI — sharp on screen and print
        import io
        img_bytes = io.BytesIO(pix.tobytes("png"))
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_bytes, 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)

    doc.close()
    out_path = os.path.join(output_dir, f"{Path(input_path).stem}.pptx")
    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# PDF ↔ Images
# ---------------------------------------------------------------------------

def pdf_to_images(input_path: str, output_dir: str, fmt: str = "png", dpi: int = 150) -> list[str]:
    import fitz
    doc = fitz.open(input_path)
    paths = []
    for i, page in enumerate(doc):
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        out = os.path.join(output_dir, f"page_{i + 1}.{fmt.lower()}")
        if fmt.lower() in ("jpg", "jpeg"):
            from PIL import Image as PILImage
            import io
            img = PILImage.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
            img.save(out, "JPEG", quality=95, optimize=True)
        else:
            pix.save(out)
        paths.append(out)
    doc.close()
    return paths


def images_to_pdf(input_paths: list[str], output_path: str) -> None:
    """Convert images → PDF. Handles any format Pillow can open (PNG/JPEG/WebP/BMP/TIFF…)."""
    import fitz
    import io
    from PIL import Image as PILImage

    if not input_paths:
        raise HTTPException(status_code=400, detail="No images provided")

    doc = fitz.open()
    for img_path in input_paths:
        with PILImage.open(img_path) as pil_img:
            # Normalise: RGBA/P → RGB, keep L (greyscale)
            if pil_img.mode not in ("RGB", "L"):
                pil_img = pil_img.convert("RGB")
            w, h = pil_img.size
            buf = io.BytesIO()
            pil_img.save(buf, format="PNG")
            img_bytes = buf.getvalue()
        page = doc.new_page(width=w, height=h)
        page.insert_image(fitz.Rect(0, 0, w, h), stream=img_bytes)

    doc.save(output_path, deflate=True)
    doc.close()


# ---------------------------------------------------------------------------
# HTML ↔ PDF  (xhtml2pdf — pure Python, better CSS than LibreOffice)
# ---------------------------------------------------------------------------

def html_to_pdf(html_content: str, output_path: str) -> None:
    """
    Convert HTML → PDF.
    Primary:  xhtml2pdf (pure Python, good CSS2 support, no system deps)
    Fallback: LibreOffice
    """
    try:
        from xhtml2pdf import pisa
        with open(output_path, "wb") as f:
            result = pisa.CreatePDF(html_content.encode("utf-8"), dest=f)
        if result.err:
            raise RuntimeError(f"xhtml2pdf error: {result.err}")
    except Exception as e:
        logger.warning("xhtml2pdf failed (%s), falling back to LibreOffice", e)
        _html_to_pdf_libreoffice(html_content, output_path)


def _html_to_pdf_libreoffice(html_content: str, output_path: str) -> None:
    temp_html = os.path.join(os.path.dirname(output_path), "_source.html")
    with open(temp_html, "w", encoding="utf-8") as f:
        f.write(html_content)
    try:
        lo_out = _libreoffice_convert(temp_html, os.path.dirname(output_path), "pdf")
        if os.path.exists(lo_out) and lo_out != output_path:
            os.replace(lo_out, output_path)
    finally:
        if os.path.exists(temp_html):
            os.remove(temp_html)


def pdf_to_html(input_path: str, output_path: str) -> None:
    """
    Convert PDF → HTML using PyMuPDF's structured HTML extraction.
    Preserves fonts, sizes, bold/italic, colours, and page layout far better
    than LibreOffice's export.
    """
    import fitz

    doc = fitz.open(input_path)
    pages_html: list[str] = []
    for i, page in enumerate(doc):
        # get_text("html") gives span-level HTML with font/size/colour attributes
        pages_html.append(
            f'<div class="page" id="page-{i + 1}" '
            f'style="position:relative;width:{page.rect.width:.0f}pt;'
            f'min-height:{page.rect.height:.0f}pt;">'
            f'{page.get_text("html")}</div>'
        )
    doc.close()

    html = (
        "<!DOCTYPE html>\n<html>\n<head>\n"
        '<meta charset="utf-8">\n'
        "<style>\n"
        "  body{font-family:Arial,sans-serif;margin:0;padding:20px;background:#e5e7eb;}\n"
        "  .page{background:#fff;margin:0 auto 24px;padding:40px;max-width:860px;"
        "border-radius:4px;box-shadow:0 2px 10px rgba(0,0,0,.15);overflow:hidden;}\n"
        "</style>\n</head>\n<body>\n"
        + "\n".join(pages_html)
        + "\n</body>\n</html>"
    )

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html)


# ---------------------------------------------------------------------------
# Markdown → PDF  (proper styled HTML → xhtml2pdf)
# ---------------------------------------------------------------------------

def markdown_to_pdf(md_text: str, output_path: str) -> None:
    import markdown as md_lib

    body = md_lib.markdown(
        md_text,
        extensions=["tables", "fenced_code", "toc", "nl2br", "sane_lists"],
    )

    html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
  @page {{ margin: 2cm; }}
  body {{ font-family: Georgia, serif; font-size: 12pt; line-height: 1.7;
          color: #1e293b; }}
  h1 {{ font-size: 22pt; font-family: Arial, sans-serif; color: #0f172a;
        border-bottom: 2px solid #2563EB; padding-bottom: 6px; margin-top: 1.4em; }}
  h2 {{ font-size: 16pt; font-family: Arial, sans-serif; color: #0f172a;
        border-bottom: 1px solid #e2e8f0; padding-bottom: 4px; margin-top: 1.2em; }}
  h3 {{ font-size: 13pt; font-family: Arial, sans-serif; color: #334155; }}
  code {{ background: #f1f5f9; padding: 2px 5px; border-radius: 3px;
          font-family: Courier, monospace; font-size: 10pt; }}
  pre  {{ background: #f1f5f9; padding: 14px; border-radius: 6px; }}
  pre code {{ background: none; padding: 0; }}
  table {{ border-collapse: collapse; width: 100%; margin: 1em 0; font-size: 10pt; }}
  th {{ background: #2563EB; color: white; padding: 8px 12px; text-align: left; }}
  td {{ border: 1px solid #e2e8f0; padding: 7px 12px; }}
  tr:nth-child(even) td {{ background: #f8fafc; }}
  blockquote {{ border-left: 4px solid #2563EB; margin: 1em 0;
                padding: 8px 16px; background: #eff6ff; color: #1e40af; }}
  ul, ol {{ padding-left: 1.5em; }}
  a {{ color: #2563EB; }}
</style>
</head>
<body>{body}</body>
</html>"""

    html_to_pdf(html, output_path)


# ---------------------------------------------------------------------------
# SVG → PNG
# ---------------------------------------------------------------------------

def svg_to_png(input_path: str, output_path: str) -> None:
    lo_out = _libreoffice_convert(input_path, os.path.dirname(output_path), "png")
    if lo_out != output_path and os.path.exists(lo_out):
        os.replace(lo_out, output_path)


# ---------------------------------------------------------------------------
# Text → PDF  (ReportLab Platypus — proper word-wrap, Unicode, page breaks)
# ---------------------------------------------------------------------------

def text_to_pdf(text: str, output_path: str, font_size: int = 12) -> None:
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors

    body_style = ParagraphStyle(
        "body",
        fontName="Helvetica",
        fontSize=font_size,
        leading=font_size * 1.45,
        textColor=colors.HexColor("#1e293b"),
        spaceAfter=3,
    )

    doc = SimpleDocTemplate(
        output_path, pagesize=A4,
        leftMargin=50, rightMargin=50,
        topMargin=50, bottomMargin=50,
    )

    story = []
    for line in text.split("\n"):
        if line.strip():
            # Escape special HTML chars so ReportLab doesn't misparse them
            safe = (line.replace("&", "&amp;")
                        .replace("<", "&lt;")
                        .replace(">", "&gt;"))
            story.append(Paragraph(safe, body_style))
        else:
            story.append(Spacer(1, font_size))

    if not story:
        story.append(Paragraph("(empty)", body_style))

    doc.build(story)


# ---------------------------------------------------------------------------
# CSV ↔ PDF / Excel
# ---------------------------------------------------------------------------

def csv_to_pdf(input_path: str, output_path: str) -> None:
    """
    CSV → PDF table via ReportLab.
    Auto-sizes columns, uses landscape for wide data, repeats header row on each page.
    """
    import csv
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import mm

    with open(input_path, newline="", encoding="utf-8-sig") as f:
        rows = list(csv.reader(f))

    if not rows:
        text_to_pdf("(Empty CSV file.)", output_path)
        return

    col_count = max(len(r) for r in rows)
    # Pad short rows
    rows = [r + [""] * (col_count - len(r)) for r in rows]

    page_size = landscape(A4) if col_count > 6 else A4
    usable_w  = page_size[0] - 30 * mm

    # Column width proportional to longest value (capped)
    raw_widths = [
        max(len(str(rows[r][c])) for r in range(min(30, len(rows)))) or 4
        for c in range(col_count)
    ]
    total = sum(raw_widths)
    col_widths = [usable_w * w / total for w in raw_widths]

    style = TableStyle([
        ("BACKGROUND",     (0, 0), (-1,  0),  colors.HexColor("#2563EB")),
        ("TEXTCOLOR",      (0, 0), (-1,  0),  colors.white),
        ("FONTNAME",       (0, 0), (-1,  0),  "Helvetica-Bold"),
        ("FONTSIZE",       (0, 0), (-1,  0),  10),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1),  [colors.white, colors.HexColor("#f1f5f9")]),
        ("FONTSIZE",       (0, 1), (-1, -1),  9),
        ("GRID",           (0, 0), (-1, -1),  0.25, colors.HexColor("#e2e8f0")),
        ("TOPPADDING",     (0, 0), (-1, -1),  6),
        ("BOTTOMPADDING",  (0, 0), (-1, -1),  6),
        ("LEFTPADDING",    (0, 0), (-1, -1),  8),
        ("RIGHTPADDING",   (0, 0), (-1, -1),  8),
        ("VALIGN",         (0, 0), (-1, -1),  "MIDDLE"),
    ])

    doc = SimpleDocTemplate(
        output_path, pagesize=page_size,
        leftMargin=15 * mm, rightMargin=15 * mm,
        topMargin=15 * mm, bottomMargin=15 * mm,
    )
    table = Table(rows, colWidths=col_widths, repeatRows=1)
    table.setStyle(style)
    doc.build([table])


def csv_to_excel(input_path: str, output_path: str) -> None:
    """
    CSV → XLSX with:
    - Blue styled header row + freeze pane
    - Auto type coercion (int / float / string)
    - Auto column widths
    - Alternating row shading
    """
    import csv
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    BLUE   = "2563EB"
    STRIPE = "F1F5F9"
    BD     = "E2E8F0"

    def _border():
        s = Side(style="thin", color=BD)
        return Border(left=s, right=s, top=s, bottom=s)

    def _coerce(v: str):
        if v == "":
            return ""
        try:
            return int(v)
        except ValueError:
            pass
        try:
            return float(v)
        except ValueError:
            pass
        return v

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"

    with open(input_path, newline="", encoding="utf-8-sig") as f:
        for r_idx, row in enumerate(csv.reader(f), 1):
            for c_idx, val in enumerate(row, 1):
                cell = ws.cell(row=r_idx, column=c_idx)
                cell.value = val if r_idx == 1 else _coerce(val)
                cell.border = _border()
                cell.alignment = Alignment(vertical="center")
                if r_idx == 1:
                    cell.fill = PatternFill("solid", fgColor=BLUE)
                    cell.font = Font(bold=True, color="FFFFFF", size=11)
                elif r_idx % 2 == 0:
                    cell.fill = PatternFill("solid", fgColor=STRIPE)

    # Auto-fit columns
    for col in ws.columns:
        max_len = max((len(str(c.value or "")) for c in col), default=8)
        ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 60)

    ws.freeze_panes = "A2"
    wb.save(output_path)


def excel_to_csv(input_path: str, output_path: str) -> None:
    import csv
    from openpyxl import load_workbook
    wb = load_workbook(input_path, read_only=True, data_only=True)
    ws = wb.active
    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for row in ws.iter_rows(values_only=True):
            writer.writerow([("" if cell is None else cell) for cell in row])


# ---------------------------------------------------------------------------
# JSON → CSV / Excel
# ---------------------------------------------------------------------------

def json_to_table(input_path: str, output_path: str, target: str = "csv") -> None:
    import json
    import csv
    from openpyxl import Workbook

    with open(input_path, encoding="utf-8") as f:
        data = json.load(f)

    # Normalise: wrap scalar/dict in a list
    if isinstance(data, dict):
        data = [data]
    elif not isinstance(data, list):
        data = [{"value": data}]

    # Flatten one level of nesting
    flat: list[dict] = []
    for item in data:
        if isinstance(item, dict):
            flat.append({k: (str(v) if isinstance(v, (dict, list)) else v)
                         for k, v in item.items()})
        else:
            flat.append({"value": item})

    keys = list(flat[0].keys()) if flat else []

    if target == "csv":
        with open(output_path, "w", newline="", encoding="utf-8") as f:
            w = csv.DictWriter(f, fieldnames=keys)
            w.writeheader()
            w.writerows(flat)
    else:
        from openpyxl.styles import Font, PatternFill
        wb = Workbook()
        ws = wb.active
        ws.append(keys)
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="2563EB")
        for row in flat:
            ws.append([row.get(k, "") for k in keys])
        for col in ws.columns:
            from openpyxl.utils import get_column_letter
            max_len = max((len(str(c.value or "")) for c in col), default=8)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 60)
        wb.save(output_path)


# ---------------------------------------------------------------------------
# EPUB → PDF
# ---------------------------------------------------------------------------

def epub_to_pdf(input_path: str, output_dir: str) -> str:
    return _libreoffice_convert(input_path, output_dir, "pdf")
