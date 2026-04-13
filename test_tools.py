"""
DocForge API Tool Tests
Runs against http://localhost:8000 with real files.
Usage: python test_tools.py
"""
import io
import os
import sys

# Force UTF-8 output on Windows
os.environ.setdefault("PYTHONIOENCODING", "utf-8")
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")


import json
import time
import fitz  # PyMuPDF
import requests
from PIL import Image

BASE = "http://localhost:8000/api/v1"

# ── Colour helpers ─────────────────────────────────────────────────────────────
GREEN = "\033[92m"
RED   = "\033[91m"
YELLOW= "\033[93m"
RESET = "\033[0m"
BOLD  = "\033[1m"

def ok(msg):   print(f"  {GREEN}[PASS]{RESET}  {msg}")
def fail(msg): print(f"  {RED}[FAIL]{RESET}  {msg}")
def skip(msg): print(f"  {YELLOW}[SKIP]{RESET}  {msg}")

# ── Sample file generators ─────────────────────────────────────────────────────

def make_pdf(pages=3, text_prefix="Test Page") -> bytes:
    """Create a simple multi-page PDF in memory."""
    doc = fitz.open()
    for i in range(pages):
        page = doc.new_page(width=595, height=842)  # A4
        page.insert_text((72, 100), f"{text_prefix} {i+1}", fontsize=24)
        page.insert_text((72, 150), "Hello DocForge! This is sample content for testing.", fontsize=12)
        page.insert_text((72, 200), f"Lorem ipsum dolor sit amet, page {i+1} of {pages}.", fontsize=12)
        # Add a simple table-like block
        for row in range(5):
            page.insert_text((72, 250 + row*20), f"Row {row+1}:  Col A  |  Col B  |  Col C", fontsize=10)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_image(fmt="JPEG") -> bytes:
    """Create a simple test image in memory."""
    img = Image.new("RGB", (400, 300), color=(73, 109, 200))
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return buf.getvalue()


def make_two_pdfs():
    return make_pdf(2, "Doc A Page"), make_pdf(2, "Doc B Page")


# ── Test runner ────────────────────────────────────────────────────────────────

results = {"pass": 0, "fail": 0, "skip": 0}

def test(name, fn):
    print(f"\n{BOLD}[{name}]{RESET}")
    t0 = time.time()
    try:
        fn()
        results["pass"] += 1
    except AssertionError as e:
        fail(str(e))
        results["fail"] += 1
    except Exception as e:
        fail(f"Exception: {e}")
        results["fail"] += 1
    print(f"  time: {time.time()-t0:.2f}s")


def assert_pdf(r, label=""):
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    ct = r.headers.get("content-type", "")
    assert "pdf" in ct, f"Expected PDF content-type, got: {ct}"
    assert len(r.content) > 500, f"Response body too small ({len(r.content)} bytes)"
    size_kb = len(r.content) / 1024
    ok(f"{label} -> {size_kb:.1f} KB PDF")


def assert_file(r, expected_ext, label=""):
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    assert len(r.content) > 100, f"Response body too small ({len(r.content)} bytes)"
    size_kb = len(r.content) / 1024
    ok(f"{label} -> {size_kb:.1f} KB ({expected_ext})")


# ═══════════════════════════════════════════════════════════════════════════════
# PDF TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def test_merge():
    pdf1, pdf2 = make_two_pdfs()
    r = requests.post(f"{BASE}/pdf/merge", files=[
        ("files", ("a.pdf", pdf1, "application/pdf")),
        ("files", ("b.pdf", pdf2, "application/pdf")),
    ])
    assert_pdf(r, "Merge 2 PDFs")


def test_split():
    pdf = make_pdf(3)
    r = requests.post(f"{BASE}/pdf/split", files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    assert "zip" in r.headers.get("content-type", ""), f"Expected ZIP, got: {r.headers.get('content-type')}"
    ok(f"Split 3-page PDF → {len(r.content)/1024:.1f} KB ZIP")


def test_compress():
    pdf = make_pdf(5)
    orig_size = len(pdf)
    r = requests.post(f"{BASE}/pdf/compress", files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert_pdf(r, f"Compress (orig={orig_size//1024}KB, result={len(r.content)//1024}KB)")


def test_rotate():
    pdf = make_pdf(2)
    r = requests.post(f"{BASE}/pdf/rotate",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"angle": "90"})
    assert_pdf(r, "Rotate 90 degrees")


def test_delete_pages():
    pdf = make_pdf(3)
    r = requests.post(f"{BASE}/pdf/delete-pages",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"page_numbers": "2"})
    assert_pdf(r, "Delete page 2 from 3-page PDF")


def test_reorder_pages():
    pdf = make_pdf(3)
    r = requests.post(f"{BASE}/pdf/reorder",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"order": "3,1,2"})
    assert_pdf(r, "Reorder pages 3,1,2")


def test_extract_pages():
    pdf = make_pdf(4)
    r = requests.post(f"{BASE}/pdf/extract-pages",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"page_numbers": "1,3"})
    assert_pdf(r, "Extract pages 1 and 3")


def test_watermark():
    pdf = make_pdf(2)
    r = requests.post(f"{BASE}/pdf/watermark",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"text": "CONFIDENTIAL", "opacity": "0.3", "font_size": "50"})
    assert_pdf(r, "Add text watermark")


def test_page_numbers():
    pdf = make_pdf(3)
    r = requests.post(f"{BASE}/pdf/page-numbers",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"position": "bottom-center", "font_size": "12", "start_number": "1"})
    assert_pdf(r, "Add page numbers")


def test_header_footer():
    pdf = make_pdf(2)
    r = requests.post(f"{BASE}/pdf/header-footer",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"header": "My Header", "footer": "My Footer", "font_size": "10"})
    assert_pdf(r, "Add header and footer")


def test_flatten():
    pdf = make_pdf(2)
    r = requests.post(f"{BASE}/pdf/flatten",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert_pdf(r, "Flatten PDF")


def test_repair():
    pdf = make_pdf(2)
    r = requests.post(f"{BASE}/pdf/repair",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert_pdf(r, "Repair PDF")


def test_metadata_get():
    pdf = make_pdf(1)
    r = requests.post(f"{BASE}/pdf/metadata",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    data = r.json()
    assert "page_count" in data, f"Missing page_count in response: {data}"
    ok(f"Metadata: {data}")


def test_encrypt_decrypt():
    pdf = make_pdf(1)
    # Encrypt — endpoint is /security/encrypt, field is user_password
    r = requests.post(f"{BASE}/security/encrypt",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"user_password": "secret123", "owner_password": ""})
    assert_pdf(r, "Encrypt with password")
    encrypted = r.content
    # Decrypt — endpoint is /security/decrypt
    r2 = requests.post(f"{BASE}/security/decrypt",
                       files=[("file", ("encrypted.pdf", encrypted, "application/pdf"))],
                       data={"password": "secret123"})
    assert_pdf(r2, "Decrypt back")


def test_redact():
    pdf = make_pdf(1)
    r = requests.post(f"{BASE}/pdf/redact",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"patterns": json.dumps(["Test Page", "Hello"])})
    assert_pdf(r, "Redact text patterns")


def test_pdf_info():
    pdf = make_pdf(3)
    r = requests.get(f"{BASE}/pdf/info", params={"page_count": True})
    # info is usually a POST with file
    r = requests.post(f"{BASE}/pdf/metadata",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert r.status_code == 200
    info = r.json()
    assert info["page_count"] == 3, f"Expected 3 pages, got {info['page_count']}"
    ok(f"PDF info: {info['page_count']} pages, {info['file_size_bytes']} bytes")


# ═══════════════════════════════════════════════════════════════════════════════
# CONVERT TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def test_pdf_to_word():
    pdf = make_pdf(2)
    r = requests.post(f"{BASE}/convert/pdf-to-word",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert_file(r, ".docx", "PDF to Word")


def test_pdf_to_excel():
    pdf = make_pdf(1)
    r = requests.post(f"{BASE}/convert/pdf-to-excel",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert_file(r, ".xlsx", "PDF to Excel")


def test_pdf_to_pptx():
    pdf = make_pdf(1)
    r = requests.post(f"{BASE}/convert/pdf-to-pptx",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))])
    assert_file(r, ".pptx", "PDF to PPTX")


def test_pdf_to_image():
    pdf = make_pdf(2)
    r = requests.post(f"{BASE}/convert/pdf-to-image",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"fmt": "png", "dpi": "72"})
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    assert "zip" in r.headers.get("content-type", ""), f"Expected ZIP"
    ok(f"PDF to Images ZIP ({len(r.content)//1024} KB)")


def test_image_to_pdf():
    img = make_image("PNG")
    # Backend: /convert/image-to-pdf uses field 'files' (plural)
    r = requests.post(f"{BASE}/convert/image-to-pdf",
                      files=[("files", ("test.png", img, "image/png"))])
    assert_pdf(r, "Image to PDF")


def test_text_to_pdf():
    r = requests.post(f"{BASE}/convert/text-to-pdf",
                      data={"text": "Hello World\nLine 2\nLine 3", "font_size": "14"})
    assert_pdf(r, "Text to PDF")


def test_csv_to_pdf():
    csv_data = b"Name,Age,City\nAlice,30,London\nBob,25,NYC\nCharlie,35,Lagos"
    r = requests.post(f"{BASE}/convert/csv-to-pdf",
                      files=[("file", ("data.csv", csv_data, "text/csv"))])
    assert_pdf(r, "CSV to PDF")


def test_csv_to_excel():
    csv_data = b"Name,Age,City\nAlice,30,London\nBob,25,NYC"
    r = requests.post(f"{BASE}/convert/csv-to-excel",
                      files=[("file", ("data.csv", csv_data, "text/csv"))])
    assert_file(r, ".xlsx", "CSV to Excel")


def test_excel_to_csv():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Age", "City"])
    ws.append(["Alice", 30, "London"])
    ws.append(["Bob", 25, "NYC"])
    buf = io.BytesIO()
    wb.save(buf)
    r = requests.post(f"{BASE}/convert/excel-to-csv",
                      files=[("file", ("data.xlsx", buf.getvalue(),
                              "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"))])
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    assert len(r.content) > 10, f"Response body too small ({len(r.content)} bytes)"
    ok(f"Excel to CSV -> {len(r.content)} bytes")


def test_word_to_pdf():
    from docx import Document
    doc = Document()
    doc.add_heading("DocForge Test", 0)
    doc.add_paragraph("This is a test Word document.")
    buf = io.BytesIO()
    doc.save(buf)
    r = requests.post(f"{BASE}/convert/word-to-pdf",
                      files=[("file", ("test.docx", buf.getvalue(),
                              "application/vnd.openxmlformats-officedocument.wordprocessingml.document"))])
    assert_pdf(r, "Word to PDF")


def test_excel_to_pdf():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Test", "Data"])
    ws.append([1, 2])
    buf = io.BytesIO()
    wb.save(buf)
    r = requests.post(f"{BASE}/convert/excel-to-pdf",
                      files=[("file", ("test.xlsx", buf.getvalue(),
                              "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"))])
    assert_pdf(r, "Excel to PDF")


def test_pptx_to_pdf():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "DocForge Test"
    buf = io.BytesIO()
    prs.save(buf)
    r = requests.post(f"{BASE}/convert/pptx-to-pdf",
                      files=[("file", ("test.pptx", buf.getvalue(),
                              "application/vnd.openxmlformats-officedocument.presentationml.presentation"))])
    assert_pdf(r, "PowerPoint to PDF")


# ═══════════════════════════════════════════════════════════════════════════════
# IMAGE TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def test_compress_image():
    img = make_image("JPEG")
    r = requests.post(f"{BASE}/image/compress",
                      files=[("file", ("photo.jpg", img, "image/jpeg"))],
                      data={"quality": "60"})
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    assert len(r.content) > 0, "Empty response"
    ok(f"Compress Image -> {len(r.content)//1024} KB ({len(r.content)} bytes)")


def test_resize_image():
    img = make_image("PNG")
    r = requests.post(f"{BASE}/image/resize",
                      files=[("file", ("photo.png", img, "image/png"))],
                      data={"width": "200", "height": "150"})
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    assert len(r.content) > 0, "Empty response"
    ok(f"Resize Image -> {len(r.content)//1024} KB ({len(r.content)} bytes)")


def test_convert_image():
    img = make_image("JPEG")
    r = requests.post(f"{BASE}/image/convert",
                      files=[("file", ("photo.jpg", img, "image/jpeg"))],
                      data={"target_format": "png"})  # field is target_format, not format
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    ok(f"Convert JPEG to PNG ({len(r.content)//1024} KB)")


# ═══════════════════════════════════════════════════════════════════════════════
# OCR TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def test_ocr_pdf():
    pdf = make_pdf(1)
    r = requests.post(f"{BASE}/ocr/pdf",
                      files=[("file", ("test.pdf", pdf, "application/pdf"))],
                      data={"lang": "eng"})
    if r.status_code == 501:
        skip("OCR PDF - Tesseract not installed (install it and set TESSERACT_PATH in .env)")
        results["skip"] += 1
        results["fail"] -= 1  # don't count as fail
        return
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:300]}"
    assert len(r.content) > 100, f"Response too small: {len(r.content)} bytes"
    ok(f"OCR PDF -> {len(r.content)//1024} KB searchable PDF")


def test_html_to_pdf():
    html = "<h1>DocForge Test</h1><p>This is a test PDF from HTML.</p>"
    r = requests.post(f"{BASE}/convert/html-to-pdf",
                      data={"html": html})
    if r.status_code == 501:
        skip("HTML to PDF - weasyprint not installed")
        results["skip"] += 1
        return
    assert_pdf(r, "HTML to PDF")


def test_svg_convert():
    svg = '<svg height="100" width="100"><circle cx="50" cy="50" r="40" stroke="black" stroke-width="3" fill="red" /></svg>'
    r = requests.post(f"{BASE}/convert/svg-convert",
                      files=[("file", ("test.svg", svg, "image/svg+xml"))],
                      data={"target": "png"})
    if r.status_code == 501:
        skip("SVG Convert - cairosvg not installed")
        results["skip"] += 1
        return
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:300]}"
    ok(f"SVG to PNG ({len(r.content)} bytes)")


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print(f"\n{BOLD}{'='*60}")
    print("  DocForge API Test Suite")
    print(f"  Target: {BASE}")
    print(f"{'='*60}{RESET}\n")

    # PDF Tools
    print(f"{BOLD}--- PDF Tools ---{RESET}")
    test("Merge PDF",           test_merge)
    test("Split PDF",           test_split)
    test("Compress PDF",        test_compress)
    test("Rotate PDF",          test_rotate)
    test("Delete Pages",        test_delete_pages)
    test("Reorder Pages",       test_reorder_pages)
    test("Extract Pages",       test_extract_pages)
    test("Watermark PDF",       test_watermark)
    test("Page Numbers",        test_page_numbers)
    test("Header and Footer",   test_header_footer)
    test("Flatten PDF",         test_flatten)
    test("Repair PDF",          test_repair)
    test("PDF Metadata",        test_metadata_get)
    test("Encrypt + Decrypt",   test_encrypt_decrypt)
    test("Redact PDF",          test_redact)
    test("PDF Info (pages)",    test_pdf_info)

    # Convert Tools
    print(f"\n{BOLD}--- Convert Tools ---{RESET}")
    test("PDF to Word",         test_pdf_to_word)
    test("Word to PDF",         test_word_to_pdf)
    test("PDF to Excel",        test_pdf_to_excel)
    test("Excel to PDF",        test_excel_to_pdf)
    test("PDF to Image (PNG)",  test_pdf_to_image)
    test("Image to PDF",        test_image_to_pdf)
    test("PDF to PPTX",         test_pdf_to_pptx)
    test("PPTX to PDF",         test_pptx_to_pdf)
    test("Text to PDF",         test_text_to_pdf)
    test("CSV to PDF",          test_csv_to_pdf)
    test("CSV to Excel",         test_csv_to_excel)
    test("Excel to CSV",         test_excel_to_csv)
    test("HTML to PDF",         test_html_to_pdf)
    test("SVG Convert",         test_svg_convert)

    # Image Tools
    print(f"\n{BOLD}--- Image Tools ---{RESET}")
    test("Compress Image",      test_compress_image)
    test("Resize Image",        test_resize_image)
    test("Convert Image",       test_convert_image)

    # OCR
    print(f"\n{BOLD}--- OCR Tools ---{RESET}")
    test("OCR PDF",             test_ocr_pdf)

    # Summary
    total = results["pass"] + results["fail"] + results["skip"]
    print(f"\n{BOLD}{'='*60}")
    print(f"  Results:  {GREEN}{results['pass']} passed{RESET}  |  {RED}{results['fail']} failed{RESET}  |  {YELLOW}{results['skip']} skipped{RESET}  /  {total} total")
    print(f"{'='*60}{RESET}\n")

    sys.exit(0 if results["fail"] == 0 else 1)
